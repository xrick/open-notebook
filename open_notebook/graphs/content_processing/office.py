from docx import Document
from loguru import logger
from openpyxl import load_workbook
from pptx import Presentation

from open_notebook.graphs.content_processing.state import SourceState

SUPPORTED_OFFICE_TYPES = [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
]


def extract_docx_content_detailed(file_path):
    try:
        doc = Document(file_path)
        content = []

        for paragraph in doc.paragraphs:
            if not paragraph.text.strip():
                continue

            style = paragraph.style.name if paragraph.style else "Normal"
            text = paragraph.text.strip()

            # Get paragraph formatting
            p_format = paragraph.paragraph_format
            indent = p_format.left_indent or 0

            # Convert indent to spaces (1 level = 4 spaces)
            indent_level = 0
            if hasattr(indent, "pt"):
                indent_level = int(indent.pt / 72)  # 72 points = 1 inch
            indent_spaces = " " * (indent_level * 4)

            # Handle different types of formatting
            if "Heading" in style:
                level = style[-1] if style[-1].isdigit() else "1"
                heading_marks = "#" * int(level)
                content.append(f"\n{heading_marks} {text}\n")

            # Handle bullet points
            elif (
                paragraph.style
                and hasattr(paragraph.style, "name")
                and paragraph.style.name.startswith("List")
            ):
                # Numbered list
                if (
                    hasattr(paragraph._p, "pPr")
                    and paragraph._p.pPr is not None
                    and hasattr(paragraph._p.pPr, "numPr")
                    and paragraph._p.pPr.numPr is not None
                ):
                    # Try to get the actual number
                    try:
                        if (
                            hasattr(paragraph._p.pPr.numPr, "numId")
                            and paragraph._p.pPr.numPr.numId is not None
                            and hasattr(paragraph._p.pPr.numPr.numId, "val")
                        ):
                            number = paragraph._p.pPr.numPr.numId.val
                            content.append(f"{indent_spaces}{number}. {text}")
                        else:
                            content.append(f"{indent_spaces}1. {text}")
                    except Exception:
                        content.append(f"{indent_spaces}1. {text}")
                # Bullet list
                else:
                    content.append(f"{indent_spaces}* {text}")

            else:
                # Handle text formatting
                formatted_text = []
                for run in paragraph.runs:
                    if run.bold:
                        formatted_text.append(f"**{run.text}**")
                    elif run.italic:
                        formatted_text.append(f"*{run.text}*")
                    else:
                        formatted_text.append(run.text)

                content.append(f"{indent_spaces}{''.join(formatted_text)}")

        return "\n\n".join(content)

    except Exception as e:
        logger.error(f"Failed to extract DOCX content: {e}")
        return None


# Example of usage with metadata
def get_docx_info(file_path):
    try:
        doc = Document(file_path)

        # Extract core properties if available
        core_props = {
            "author": doc.core_properties.author,
            "created": doc.core_properties.created,
            "modified": doc.core_properties.modified,
            "title": doc.core_properties.title,
            "subject": doc.core_properties.subject,
            "keywords": doc.core_properties.keywords,
            "category": doc.core_properties.category,
            "comments": doc.core_properties.comments,
        }

        # Get document content
        content = extract_docx_content_detailed(file_path)

        # Get document statistics
        stats = {
            "paragraph_count": len(doc.paragraphs),
            "word_count": sum(
                len(p.text.split()) for p in doc.paragraphs if p.text.strip()
            ),
            "character_count": sum(
                len(p.text) for p in doc.paragraphs if p.text.strip()
            ),
        }

        return {"metadata": core_props, "content": content, "statistics": stats}

    except Exception as e:
        logger.error(f"Failed to get DOCX info: {e}")
        return None


def extract_pptx_content(file_path):
    try:
        prs = Presentation(file_path)
        content = []

        for slide_number, slide in enumerate(prs.slides, 1):
            content.append(f"\n# Slide {slide_number}\n")

            # Extract title
            if slide.shapes.title:
                content.append(f"## {slide.shapes.title.text}\n")

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:  # Skip title as it's already added
                        content.append(shape.text.strip())

        return "\n\n".join(content)

    except Exception as e:
        logger.error(f"Failed to extract PPTX content: {e}")
        return None


def extract_xlsx_content(file_path, max_rows=1000, max_cols=100):
    try:
        wb = load_workbook(file_path, data_only=True)
        content = []

        for sheet in wb.sheetnames:
            ws = wb[sheet]
            content.append(f"\n# Sheet: {sheet}\n")

            # Get the maximum row and column with data
            max_row = min(ws.max_row, max_rows)
            max_col = min(ws.max_column, max_cols)

            # Create markdown table header
            headers = []
            for col in range(1, max_col + 1):
                cell_value = ws.cell(row=1, column=col).value
                headers.append(str(cell_value) if cell_value is not None else "")

            content.append("| " + " | ".join(headers) + " |")
            content.append("| " + " | ".join(["---"] * len(headers)) + " |")

            # Add table content
            for row in range(2, max_row + 1):
                row_data = []
                for col in range(1, max_col + 1):
                    cell_value = ws.cell(row=row, column=col).value
                    row_data.append(str(cell_value) if cell_value is not None else "")
                content.append("| " + " | ".join(row_data) + " |")

        return "\n".join(content)

    except Exception as e:
        logger.error(f"Failed to extract XLSX content: {e}")
        return None


def get_pptx_info(file_path):
    try:
        prs = Presentation(file_path)

        # Extract basic properties
        props = {
            "slide_count": len(prs.slides),
            "title": "",  # PowerPoint doesn't have built-in metadata like Word
        }

        # Get document content
        content = extract_pptx_content(file_path)

        # Get presentation statistics
        stats = {
            "slide_count": len(prs.slides),
            "shape_count": sum(len(slide.shapes) for slide in prs.slides),
            "text_frame_count": sum(
                sum(1 for shape in slide.shapes if hasattr(shape, "text"))
                for slide in prs.slides
            ),
        }

        return {"metadata": props, "content": content, "statistics": stats}

    except Exception as e:
        logger.error(f"Failed to get PPTX info: {e}")
        return None


def get_xlsx_info(file_path):
    try:
        wb = load_workbook(file_path, data_only=True)

        # Extract basic properties
        props = {
            "sheet_count": len(wb.sheetnames),
            "sheets": wb.sheetnames,
            "title": wb.properties.title,
            "creator": wb.properties.creator,
            "created": wb.properties.created,
            "modified": wb.properties.modified,
        }

        # Get document content
        content = extract_xlsx_content(file_path)

        # Get workbook statistics
        stats = {
            "sheet_count": len(wb.sheetnames),
            "total_rows": sum(sheet.max_row for sheet in wb.worksheets),
            "total_columns": sum(sheet.max_column for sheet in wb.worksheets),
        }

        return {"metadata": props, "content": content, "statistics": stats}

    except Exception as e:
        logger.error(f"Failed to get XLSX info: {e}")
        return None


def extract_office_content(state: SourceState):
    """Universal function to extract content from Office files"""
    assert state.get("file_path"), "No file path provided"
    assert (
        state.get("identified_type") in SUPPORTED_OFFICE_TYPES
    ), "Unsupported File Type"

    file_path = state["file_path"]
    doc_type = state["identified_type"]

    if (
        doc_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        logger.debug("Extracting content from DOCX file")
        content = extract_docx_content_detailed(file_path)
        info = get_docx_info(file_path)
    elif (
        doc_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        logger.debug("Extracting content from PPTX file")
        content = extract_pptx_content(file_path)
        info = get_pptx_info(file_path)
    elif (
        doc_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        logger.debug("Extracting content from XLSX file")
        content = extract_xlsx_content(file_path)
        info = get_xlsx_info(file_path)
    else:
        raise Exception(f"Unsupported file format: {doc_type}")

    del info["content"]

    return {"content": content, "metadata": info}
